# -*- coding: utf-8 -*-
"""
    ragnarok.document_loaders.pptx
    ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~

    Microsoft PowerPoint presentation (pptx) document loader.
"""

from pathlib import Path

from langchain_core.documents import Document
from langchain_text_splitters import TextSplitter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def unwrap_group(group):
    """
    Get a list of shapes within a ppt group (recursive).

    :param group: ppt group object (has attribute <shapes>)
    :return: shapes tuple
    """

    out = []
    for shape in group.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(unwrap_group(shape))
        else:
            out.append(shape)
    return out


def table_format(data) -> str:
    """
    Format list of lists in tabular format.

    :param data: input list of lists (rows)
    :return: string data in tabular format
    """

    # add cells if needed
    columns = max(len(row) for row in data)
    data = [row + [""] * (columns - len(row)) for row in data]

    # determine cell widths
    # noinspection PyTypeChecker
    tmp = list(map(list, zip(*data)))
    cell_widths = [max(len(x) for x in col) for col in tmp]

    # construct the row format string
    row_format = "|"
    for cw in cell_widths:
        row_format += f" {{:<{cw}}} |"

    output = [row_format.format(*row) for row in data]
    return "\n".join(output)


def table_text(table) -> str:
    """
    Retrieve text from a ppt table object and process it.

    :param table: ppt table object (property has_table == True)
    :return: processed text
    """

    data = [[cell.text.strip() for cell in row.cells] for row in table.table.rows]
    data = [["; ".join([x.strip() for x in t.split("\n") if x]) for t in row] for row in data]
    return table_format(data)


def shape_text(shape) -> str | None:
    """
    Retrieve text from a ppt shape object and process it.

    :param shape: ppt shape object
    :return: processed text
    """

    if shape.has_table:
        return table_text(shape)

    if not shape.has_text_frame:
        return None

    text = [p.text.strip() for p in shape.text_frame.paragraphs]
    text = [" ".join(p.split()) for p in text if p]
    return "\n".join(text)


class PyPPTXLoader:

    def __init__(self, file_path: str | Path):
        self.file_path = file_path

    def load(self) -> list[Document]:

        documents = []
        prs = Presentation(self.file_path)

        for page, slide in enumerate(prs.slides):
            shapes = unwrap_group(slide)
            content = [shape_text(shape) for shape in shapes]
            title = [", ".join(t.split("\n")) for t, s in zip(content, shapes) if "title" in s.name.lower()]
            title = "\n".join(title)
            content = [text for text in content if text]
            content = "\n".join(content)

            if len(content) < 50:
                continue

            documents.append(Document(page_content=content, metadata={"page": page, "title": title}))

        return documents

    def load_and_split(self, text_splitter: TextSplitter) -> list[Document]:
        return text_splitter.split_documents(self.load())
